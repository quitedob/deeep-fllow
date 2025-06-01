# 文件路径: src/tools/output_generator.py
# -*- coding: utf-8 -*-
"""
OutputGenerator：将 Markdown 内容导出为 TXT、PDF、PPTX，并返回对应路径。
"""
import os
import aiofiles
import markdown2
import pdfkit # type: ignore
from pptx import Presentation
from typing import Dict, Any
import logging

logger = logging.getLogger(__name__)

class OutputGenerator:
    def __init__(self, config: Dict[str, Any]):
        """
        config:
          - output_dir: 输出目录（相对或绝对路径）
        """
        self.output_dir = config.get("output_dir", "outputs")
        try:
            os.makedirs(self.output_dir, exist_ok=True)
            logger.info(f"输出目录 '{self.output_dir}' 已确保存在。")
        except OSError as e:
            logger.error(f"创建输出目录 '{self.output_dir}' 失败: {e}", exc_info=True)
            raise

    async def to_text(self, md: str) -> str:
        """
        将 Markdown 转为纯文本（去掉标记），并写入 .txt 文件。
        """
        txt_path = os.path.join(self.output_dir, "report.txt")
        logger.info(f"准备将 Markdown 转换为 TXT: {txt_path}")
        try:
            # 修复：使用 aiofiles
            async with aiofiles.open(txt_path, "w", encoding="utf-8") as f:
                await f.write(md)
            logger.info(f"TXT 文件已成功保存至: {txt_path}")
            return txt_path
        except Exception as e:
            logger.error(f"保存 TXT 文件失败 ({txt_path}): {e}", exc_info=True)
            raise

    async def to_markdown(self, md: str) -> str:
        """
        直接输出 Markdown 文件（.md）。
        """
        md_path = os.path.join(self.output_dir, "report.md")
        logger.info(f"准备保存 Markdown 文件: {md_path}")
        try:
            # 修复：使用 aiofiles
            async with aiofiles.open(md_path, "w", encoding="utf-8") as f:
                await f.write(md)
            logger.info(f"Markdown 文件已成功保存至: {md_path}")
            return md_path
        except Exception as e:
            logger.error(f"保存 Markdown 文件失败 ({md_path}): {e}", exc_info=True)
            raise

    async def to_pdf(self, md: str) -> str:
        """
        将 Markdown 转为 HTML，再调用 pdfkit 生成 PDF。
        依赖：系统需安装 wkhtmltopdf。
        """
        pdf_path = os.path.join(self.output_dir, "report.pdf")
        logger.info(f"准备将 Markdown 转换为 PDF: {pdf_path}")
        try:
            html = markdown2.markdown(md, extras=["fenced-code-blocks", "tables", "footnotes", "header-ids", "smarty-pants"])
            options = {
                'encoding': "UTF-8",
                'custom-header': [
                    ('Content-Encoding', 'utf-8'),
                ],
                'no-outline': None,
                'quiet': ''
            }
            pdfkit.from_string(html, pdf_path, options=options)
            logger.info(f"PDF 文件已成功生成至: {pdf_path}")
            return pdf_path
        except FileNotFoundError:
            logger.error(f"生成 PDF 文件失败: wkhtmltopdf 未找到。请确保已安装并将其添加到系统 PATH。")
            raise RuntimeError("wkhtmltopdf not found. Please install it and ensure it's in your PATH.")
        except Exception as e:
            logger.error(f"生成 PDF 文件失败 ({pdf_path}): {e}", exc_info=True)
            if isinstance(e, OSError) and "wkhtmltopdf" in str(e):
                raise RuntimeError(f"wkhtmltopdf execution failed. Is it installed and in PATH? Original error: {e}")
            raise

    async def to_ppt(self, ppt_json: Dict[str, Any]) -> str:
        """
        将自定义的 ppt_json 转换为 PPTX。
        """
        ppt_path = os.path.join(self.output_dir, "report.pptx")
        logger.info(f"准备从 JSON 生成 PPTX: {ppt_path}")
        try:
            prs = Presentation()
            slide_layout_title = prs.slide_layouts[0]
            slide = prs.slides.add_slide(slide_layout_title)
            title_shape = slide.shapes.title
            subtitle_shape = slide.placeholders.get(1)

            title_shape.text = ppt_json.get("title", "演示文稿")
            if subtitle_shape:
                subtitle_shape.text = ""

            slide_layout_content = prs.slide_layouts[1]
            for slide_info in ppt_json.get("slides", []):
                slide = prs.slides.add_slide(slide_layout_content)
                title_shape = slide.shapes.title
                body_shape = slide.shapes.placeholders[1]

                title_shape.text = slide_info.get("heading", "内容页")

                tf = body_shape.text_frame
                tf.clear()

                content_text = slide_info.get("content", "")
                for line in content_text.split('\n'):
                    p = tf.add_paragraph()
                    p.text = line.strip()

            prs.save(ppt_path)
            logger.info(f"PPTX 文件已成功生成至: {ppt_path}")
            return ppt_path
        except Exception as e:
            logger.error(f"生成 PPTX 文件失败 ({ppt_path}): {e}", exc_info=True)
            raise