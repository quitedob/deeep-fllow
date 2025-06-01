# 文件路径: src/tools/ppt_generator.py
# -*- coding: utf-8 -*-
# 主要功能：从结构化的JSON数据生成PPTX演示文稿。

import pptx
from pptx.util import Inches, Pt
from typing import Dict, Any, List # 修复：添加 List 类型提示

# 注意：原先在此文件中重复定义的 Configuration 类已被移除。
# 如果此模块未来需要配置，应从 src.config.configuration 导入。

def generate_ppt_from_json(data: Dict[str, Any], output_path: str) -> str:
    """
    根据输入的JSON数据生成一个.pptx演示文稿。

    :param data: 包含'title'和'slides'键的字典。
                 'slides'是一个列表，每个元素是包含'title'和'points'的字典。
    :param output_path: 生成的 .pptx 文件的保存路径。
    :return: 一个表示操作状态的字符串。
    简化注释：从JSON生成PPT
    """
    try:
        # 1. 验证输入数据的基本结构
        if 'title' not in data or 'slides' not in data or not isinstance(data['slides'], list): # [cite: 640]
            raise ValueError("JSON数据缺少 'title' 或 'slides' 键，或者 'slides' 不是一个列表。")

        # 2. 创建一个新的演示文稿，使用16:9宽屏布局
        prs = pptx.Presentation()
        prs.slide_width = Inches(16) # [cite: 639]
        prs.slide_height = Inches(9) # [cite: 639]

        # 3. 添加标题页 (布局0: 标题幻灯片)
        title_slide_layout = prs.slide_layouts[0]
        slide = prs.slides.add_slide(title_slide_layout) # [cite: 641]
        title = slide.shapes.title
        subtitle = slide.placeholders.get(1)  # 安全地获取副标题占位符 # [cite: 641]

        title.text = data.get('title', '无标题报告')

        # 从第一个幻灯片内容中提取副标题，或使用默认值
        # 修复：确保在访问 slides[0] 前检查 slides 是否为空
        if data['slides'] and isinstance(data['slides'][0], dict) and isinstance(data['slides'][0].get('points'), list): # [cite: 319, 641]
            if subtitle:
                subtitle.text = "\n".join(data['slides'][0].get('points', []))
        elif subtitle: # [cite: 642]
            subtitle.text = "由 AI Agent 生成"

        # 设置字体大小
        if title.text_frame.paragraphs: # 确保有段落
            title.text_frame.paragraphs[0].font.size = Pt(44) # [cite: 642]
        if subtitle and subtitle.text_frame.paragraphs: # 确保有段落
            subtitle.text_frame.paragraphs[0].font.size = Pt(24) # [cite: 642]

        # 4. 遍历数据，添加内容幻灯片 (跳过已用作封面的第一个，如果存在)
        content_layout = prs.slide_layouts[1]  # 布局1: "标题和内容"
        slides_data_iter = iter(data['slides'])
        if data['slides']: # 跳过第一个（已用作封面部分信息）
            next(slides_data_iter, None)

        for item in slides_data_iter:
            if not isinstance(item, dict) or 'title' not in item or 'points' not in item: # [cite: 643]
                print(f"--- [PPT生成器] 警告: 跳过格式不正确的幻灯片数据项: {item} ---")
                continue

            slide = prs.slides.add_slide(content_layout) # [cite: 322]

            title_shape = slide.shapes.title
            title_shape.text = item.get('title', '无标题页面')
            if title_shape.text_frame.paragraphs:
                title_shape.text_frame.paragraphs[0].font.size = Pt(36) # [cite: 644]

            body_shape = slide.shapes.placeholders[1]
            tf = body_shape.text_frame
            tf.clear()  # 清除默认文本，防止出现"单击此处添加文本" # [cite: 644]

            points = item.get('points', []) # [cite: 323]
            # 兼容要点是单个字符串或字符串列表的情况
            if isinstance(points, list):
                if not points:  # 如果要点列表为空，添加一个提示 # [cite: 323]
                    p = tf.add_paragraph()
                    p.text = "(此页无内容)" # [cite: 645]
                    p.font.size = Pt(18) # [cite: 645]
                else:
                    for point_text in points: # 变量名修改避免与外层 point 混淆
                        p = tf.add_paragraph()
                        p.text = str(point_text) # [cite: 325, 646]
                        p.font.size = Pt(22) # [cite: 325, 647]
                        p.level = 0 # [cite: 325, 647]
            else: # 如果 points 是单个字符串
                p = tf.add_paragraph()
                p.text = str(points) # [cite: 326]
                p.font.size = Pt(22) # [cite: 326]

        # 5. 保存演示文稿
        prs.save(output_path) # [cite: 648]
        return f"PPT 文件已成功保存到: {output_path}"

    except Exception as e:
        error_message = f"生成PPT文件时发生错误: {e}"
        print(f"--- [PPT生成器] {error_message} ---")
        # 重新抛出异常，让上层调用者知道发生了错误
        raise RuntimeError(error_message) from e